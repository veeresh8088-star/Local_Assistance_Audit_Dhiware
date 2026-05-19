# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
import requests
import pdfplumber
import time, json, hashlib, uuid
from datetime import datetime
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, text
from sqlalchemy.orm import DeclarativeBase, sessionmaker
from docx import Document

st.set_page_config(page_title="AICyberAuditBox", page_icon="🛡️", layout="wide")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap');
html, body, [class*="css"] { font-family: 'Inter', sans-serif !important; }

section[data-testid="stSidebar"] .stButton>button {
    background: linear-gradient(135deg,#3b82f6,#1d4ed8); color:white !important;
    border:none; border-radius:8px; font-weight:600; transition:.2s;
}
section[data-testid="stSidebar"] .stButton>button:hover { transform:translateY(-2px); box-shadow:0 4px 15px rgba(59,130,246,.4); }
.main-header { background:linear-gradient(135deg,#0f172a 0%,#1e3a5f 100%);
    padding:28px 32px; border-radius:16px; margin-bottom:24px;
    border:1px solid rgba(59,130,246,.2); }
.stat-card { background:#1e293b; border:1px solid #334155; border-radius:12px;
    padding:20px; text-align:center; }
.stat-num { font-size:2rem; font-weight:700; }
.badge-critical { background:#450a0a; border-left:4px solid #ef4444; border-radius:8px; padding:16px; margin:8px 0; }
.badge-high     { background:#431407; border-left:4px solid #f97316; border-radius:8px; padding:16px; margin:8px 0; }
.badge-medium   { background:#422006; border-left:4px solid #eab308; border-radius:8px; padding:16px; margin:8px 0; }
.badge-low      { background:#052e16; border-left:4px solid #22c55e; border-radius:8px; padding:16px; margin:8px 0; }
.chat-bubble-user { background:#1e3a5f; border-radius:16px 16px 4px 16px; padding:12px 16px;
    margin:8px 0; margin-left:15%; color:#e2e8f0; }
.chat-bubble-bot  { background:#1e293b; border:1px solid #334155; border-radius:16px 16px 16px 4px;
    padding:12px 16px; margin:8px 0; margin-right:5%; color:#e2e8f0; }
.uc-card { background:#1e293b; border:1px solid #334155; border-radius:10px;
    padding:14px 18px; margin:8px 0; cursor:pointer; transition:.2s; }
.uc-card:hover { border-color:#3b82f6; transform:translateX(4px); }
.stage-done   { color:#22c55e; border-left:3px solid #22c55e; padding:6px 0 6px 14px; margin:4px 0; font-weight:600; }
.stage-active { color:#3b82f6; border-left:3px solid #3b82f6; padding:6px 0 6px 14px; margin:4px 0; font-weight:600; }
.stage-idle   { color:#475569; border-left:3px solid #334155; padding:6px 0 6px 14px; margin:4px 0; }
div[data-testid="stDecoration"] { display:none; }
</style>
""", unsafe_allow_html=True)

# ── USE CASES ─────────────────────────────────────────────────────────────────
USE_CASES = [
    # --- ISO 27001 ---
    {"sl":1,"standard":"ISO 27001","label":"Access Control Policy (A.5.15)","icon":"🔐","use_case":"Access Control Policy – Grant, Review, Revoke access","expected":"Verify documented access control policy covering full access lifecycle.","format":"PDF","prompt_hint":"Verify if a documented access control policy exists covering granting, reviewing, and revoking access."},
    {"sl":2,"standard":"ISO 27001","label":"Role-Based Access Control (RBAC)","icon":"👥","use_case":"RBAC – Permissions assigned to Roles, not individuals","expected":"Verify RBAC implementation. Access should be role-based.","format":"PDF","prompt_hint":"Verify if access is managed through roles rather than individual users. Identify RBAC gaps."},
    {"sl":3,"standard":"ISO 27001","label":"Multi-Factor Authentication (MFA)","icon":"🔑","use_case":"MFA – Enforced for all external access via VPN / cloud","expected":"Confirm MFA enforcement, password complexity and rotation requirements.","format":"PDF","prompt_hint":"Verify MFA enforcement for external access, VPN, cloud. Check password complexity and rotation policy."},
    {"sl":4,"standard":"ISO 27001","label":"Privileged Access Management","icon":"⚡","use_case":"Privileged Access – Time-limited and monitored","expected":"Verify privileged access is restricted, time-limited, and monitored.","format":"PDF","prompt_hint":"Verify if privileged access is time-limited, restricted to legitimate need, and under enhanced monitoring."},
    {"sl":5,"standard":"ISO 27001","label":"Access Reviews & Orphaned Accounts","icon":"🔎","use_case":"Access Reviews – Periodic review and orphaned account management","expected":"Verify periodic access reviews and prompt revocation. Orphaned accounts managed.","format":"PDF","prompt_hint":"Check if access rights are reviewed periodically, revoked promptly, and orphaned accounts are managed."},
    {"sl":6,"standard":"ISO 27001","label":"Incident Mgmt – Vendor Assessment","icon":"🔍","use_case":"Incident Management (A.5.24) – Vendor Security Assessment","expected":"Verify vendor security measures comply with ISO 27001. Identify gaps.","format":"PDF","prompt_hint":"Verify if vendor security measures comply with ISO 27001 A.5.24-A.5.28. List all gaps."},
    {"sl":7,"standard":"ISO 27001","label":"Incident Mgmt – Policy Review","icon":"🔄","use_case":"Incident Management (A.5.28) – Conduct Regular Reviews","expected":"Verify if policies are regularly reviewed and updated. Identify stale policies.","format":"PDF","prompt_hint":"Check if incident response policy is regularly reviewed. Identify outdated content."},
    
    # --- DPDP / GDPR ---
    {"sl":8,"standard":"DPDP / GDPR","label":"Consent Management & Notice","icon":"📝","use_case":"Verify Consent mechanisms and Privacy Notice transparency","expected":"Confirm clear, granular, revocable consent notice complying with DPDP/GDPR.","format":"PDF","prompt_hint":"Check privacy policy and notice for consent clarity, purpose limitation, and DPO details."},
    {"sl":9,"standard":"DPDP / GDPR","label":"Data Protection Officer (DPO)","icon":"👔","use_case":"Verify appointment of DPO and contact availability","expected":"Confirm DPO details are published and contact details accessible.","format":"PDF","prompt_hint":"Search for Data Protection Officer (DPO) designation and contact email in policies."},
    {"sl":10,"standard":"DPDP / GDPR","label":"Data Subject Rights (DSR/DSAR)","icon":"👤","use_case":"Verify procedures for handling DSAR requests","expected":"Confirm response SLA for data deletion, access, and correction requests.","format":"PDF","prompt_hint":"Verify DSAR processing SLA, data deletion, correction procedures."},

    # --- SOC 2 ---
    {"sl":11,"standard":"SOC 2","label":"Security: Firewall & Encryption","icon":"🧱","use_case":"CC6.6, CC6.7 - Encryption in transit and rest","expected":"Verify firewall controls and SSL/TLS and AES encryption enforcement.","format":"PDF","prompt_hint":"Check encryption protocols in transit (TLS 1.2+) and at rest (AES-256)."},
    {"sl":12,"standard":"SOC 2","label":"Availability: Backup & Recovery","icon":"💾","use_case":"CC7.5 - Backup restoration and disaster recovery planning","expected":"Verify automated backups and disaster recovery runbook availability.","format":"PDF","prompt_hint":"Verify daily automated backups, offsite retention, and DR testing plan."},

    # --- BCMS ---
    {"sl":13,"standard":"BCMS (Business Continuity)","label":"BCMS Continuity & ISO Certificates","icon":"🏅","use_case":"Check the Stale/Expired ISO Certificates","expected":"Validate ISO/BCMS certifications, Risk/Severity and mitigation recommendation","format":"PDF","prompt_hint":"Check if ISO/BCMS certificate is expired or expiring. Provide risk and recommendation."},
    {"sl":14,"standard":"BCMS (Business Continuity)","label":"BCP Drill & Test Results","icon":"🏃‍♂️","use_case":"Verify annual BCP testing and drill execution","expected":"Verify BCP drill results and RTO/RPO performance verification.","format":"PDF","prompt_hint":"Search for Business Continuity Plan (BCP) testing dates and results inside logs."},

    # --- X-BOM ---
    {"sl":15,"standard":"X-BOM (Software Bill of Materials)","label":"License Agreement Validity","icon":"📄","use_case":"Check and summarize the validity of the license agreement","expected":"License Type, validity date, EOL/EOS status, Risk/Severity and recommendation","format":"PDF","prompt_hint":"Summarize the license type, validity dates. Identify if EOL/EOS. Provide risk severity and recommendation."},
    {"sl":16,"standard":"X-BOM (Software Bill of Materials)","label":"Third-party Disposal (Media A.7.10)","icon":"♻️","use_case":"Third-party EWaste disposal agreement – Media Handling (A.7.10)","expected":"Verify the validity of the EWaste Agreement certificate","format":"DOC","prompt_hint":"Verify the validity date and terms of the EWaste disposal agreement. Check if current and compliant."}
]

DEMO_FINDINGS = {
    1: [{"severity":"CRITICAL","control":"ISO 27001 A.5.15","finding":"No documented access control policy found in uploaded evidence.","recommendation":"Create and publish a formal Access Control Policy covering grant/review/revoke lifecycle."}],
    2: [{"severity":"HIGH","control":"ISO 27001 A.5.15 RBAC","finding":"Access granted on individual basis. No role-based model documented.","recommendation":"Implement RBAC model and document role definitions in the policy."}],
    3: [{"severity":"CRITICAL","control":"ISO 27001 A.8.5 / NIST IA-2","finding":"No MFA policy for VPN or cloud external access found in evidence.","recommendation":"Enforce MFA for all external access. Document password complexity and 90-day rotation."}],
    4: [{"severity":"HIGH","control":"ISO 27001 A.8.2","finding":"No time-limiting or enhanced monitoring of privileged accounts documented.","recommendation":"Implement Just-In-Time (JIT) privileged access with PAM tool logging and automated expiry."}],
    5: [{"severity":"HIGH","control":"ISO 27001 A.5.18","finding":"No evidence of periodic access reviews or orphaned account removal process.","recommendation":"Implement quarterly access review process with documented approvals."}],
    6: [{"severity":"HIGH","control":"ISO 27001 A.5.24","finding":"Incident Response Plan lacks vendor-specific security assessment clauses.","recommendation":"Add vendor security assessment section aligned with ISO 27001 A.5.24–A.5.28."}],
    7: [{"severity":"MEDIUM","control":"ISO 27001 A.5.28","finding":"Policy document last reviewed in 2021. No annual review evidence found.","recommendation":"Establish a documented annual review cycle with CISO sign-off."}],
    
    # DPDP / GDPR
    8: [{"severity":"CRITICAL","control":"DPDP Sec 6 / GDPR Art 7","finding":"Privacy notice lack granular consent options. Pre-checked boxes found for marketing.","recommendation":"Implement explicit, opt-in consent and uncheck marketing boxes by default."}],
    9: [{"severity":"HIGH","control":"DPDP Sec 10 / GDPR Art 37","finding":"DPO designation details are missing from the public privacy policy document.","recommendation":"Publish Data Protection Officer name, email, and postal address in the privacy notice."}],
    10: [{"severity":"HIGH","control":"DPDP Sec 12 / GDPR Art 15","finding":"DSAR policy does not specify statutory response timeframe (30 days for GDPR).","recommendation":"Update DSAR procedure to guarantee responses within 30 days and document DSR verification process."}],
    
    # SOC 2
    11: [{"severity":"CRITICAL","control":"SOC 2 CC6.6 / CC6.7","finding":"Production data transmitted over HTTP (unencrypted) in internal API endpoints.","recommendation":"Enforce HTTPS (TLS 1.3) across all internal microservices and disable SSLv3/TLS1.0."}],
    12: [{"severity":"HIGH","control":"SOC 2 CC7.5","finding":"DR Plan is present, but recovery restoration drills have not been performed or verified in 2025.","recommendation":"Schedule and execute a mock database recovery drill, and document the actual RTO/RPO achieved."}],
    
    # BCMS
    13: [{"severity":"CRITICAL","control":"ISO 22301 Clause 9.1","finding":"ISO/BCMS Certificate expired on 2026-03-15. Certificate is no longer valid.","recommendation":"Initiate recertification audit immediately through an accredited body."}],
    14: [{"severity":"HIGH","control":"BCMS Continuity","finding":"No evidence of BCP testing or simulation drills in the last 12 months.","recommendation":"Conduct a BCP drill and document results before next audit."}],
    
    # X-BOM
    15: [{"severity":"CRITICAL","control":"Asset Mgmt / License","finding":"License expired on 2016-04-22 — 10 years ago. EOL confirmed with no vendor support.","recommendation":"Immediately replace or renew PJSIP software license to mitigate legal and security risk."}],
    16: [{"severity":"CRITICAL","control":"ISO 27001 A.7.10","finding":"EWaste Agreement Certificate is expired or not present in uploaded document.","recommendation":"Renew the third-party EWaste disposal agreement certificate immediately."}],
    
    "CROSS_FILE": [
        {"severity":"CRITICAL","control":"Cross-Document Correlation","finding":"Policy PDF (File 1) mandates 90-day password rotation, but Evidence Certificate (File 2) shows rotation set to 180 days.","recommendation":"Sync the actual system settings with the written policy document."},
        {"severity":"HIGH","control":"Cross-Document Correlation","finding":"Incident Plan (File 1) lists an external vendor for forensics, but the vendor contract (File 2) has been expired for 6 months.","recommendation":"Renew the vendor contract or update the Incident Plan with a new forensic partner."}
    ]
}

# Resolution keywords: if any of these appear in uploaded content, the gap is resolved
GAP_RESOLUTION = {
    "ISO 27001 A.5.15":            ["access control policy", "grant review revoke", "access policy document", "access control", "user access", "authorization policy"],
    "ISO 27001 A.5.15 RBAC":       ["role based", "rbac", "role assignment", "roles defined", "role-based access", "role-based"],
    "ISO 27001 A.8.5 / NIST IA-2":["mfa enabled", "multi-factor", "two-factor", "2fa", "authenticator app", "otp", "mfa", "2fa", "authenticator"],
    "ISO 27001 A.8.2":             ["privileged access", "pam tool", "just-in-time", "jit access", "time-limited access", "pam", "jit"],
    "ISO 27001 A.5.18":            ["access review completed", "quarterly review", "orphaned account removed", "account audit", "access review", "user review"],
    "ISO 27001 A.5.24":            ["vendor assessment", "vendor security", "third party assessment", "supplier review", "vendor", "third-party", "supplier"],
    "ISO 27001 A.5.28":            ["annual review", "policy reviewed 202", "reviewed and approved", "ciso sign", "annual review", "approved by ciso"],
    "DPDP Sec 6 / GDPR Art 7":     ["opt-in consent", "granular consent", "consent notice", "explicit consent", "consent form", "opt-in"],
    "DPDP Sec 10 / GDPR Art 37":   ["dpo email", "appointed dpo", "dpo details", "data protection officer", "dpo"],
    "DPDP Sec 12 / GDPR Art 15":   ["dsar SLA", "dsar response", "data subject rights", "30 days SLA", "dsar"],
    "SOC 2 CC6.6 / CC6.7":         ["tls 1.2", "tls 1.3", "https enforced", "aes-256", "encryption in transit", "tls", "https", "ssl", "encryption"],
    "SOC 2 CC7.5":                 ["dr test", "restore verify", "disaster recovery test", "backup test", "dr test", "disaster recovery", "backup"],
    "ISO 22301 Clause 9.1":        ["iso certified", "certificate valid", "certification active", "audit passed", "recertified", "iso certification", "certificate"],
    "BCMS Continuity":             ["bcp test", "drill conducted", "recovery test", "continuity test", "rto rpo", "bcp", "business continuity"],
    "Asset Mgmt / License":        ["license renewed", "new license", "valid license", "commercial agreement", "license valid until", "software license"],
    "ISO 27001 A.7.10":            ["e-waste", "ewaste", "disposal certificate", "media disposal", "certificate of destruction", "it asset disposal", "waste agreement", "ewaste", "e-waste", "disposal"],
}

# ── DATABASE ──────────────────────────────────────────────────────────────────
class Base(DeclarativeBase): pass
class AuditFinding(Base):
    __tablename__ = "audit_findings"
    id             = Column(Integer, primary_key=True, autoincrement=True)
    use_case_sl    = Column(Integer)
    use_case_name  = Column(String(300))
    severity       = Column(String(50))
    control        = Column(String(200))
    finding        = Column(Text)
    recommendation = Column(Text)
    status         = Column(String(50), default="Open")
    comment        = Column(Text, default="")
    source_files   = Column(Text, default="All uploaded documents")
    created_at     = Column(DateTime, default=datetime.utcnow)

class ChatMessage(Base):
    __tablename__ = "chat_messages"
    id             = Column(Integer, primary_key=True, autoincrement=True)
    session_id     = Column(String(100))
    session_title  = Column(String(300))
    role           = Column(String(50))
    content        = Column(Text)
    created_at     = Column(DateTime, default=datetime.utcnow)

@st.cache_resource
def init_db():
    try:
        # Try ShaktiDB (PostgreSQL)
        eng = create_engine("postgresql://postgres:ShakthiDB%402026@localhost:15234/postgres", connect_args={"connect_timeout":3})
        with eng.connect() as c: c.execute(text("SELECT 1"))
        
        # Self-healing: Check if the schema is up to date, if not, reset the table
        try:
            with eng.connect() as c: c.execute(text("SELECT source_files FROM audit_findings LIMIT 1"))
        except:
            # Table might not exist or has an old schema, drop and recreate
            from sqlalchemy import MetaData
            meta = MetaData()
            meta.reflect(bind=eng)
            if "audit_findings" in meta.tables:
                meta.tables["audit_findings"].drop(bind=eng)
        
        Base.metadata.create_all(bind=eng)
        return eng, "ShaktiDB"
    except Exception as e:
        # Fallback to Local SQLite
        eng = create_engine("sqlite:///shakthidb_local.db")
        
        # Self-healing for Local DB as well
        try:
            with eng.connect() as c: c.execute(text("SELECT source_files FROM audit_findings LIMIT 1"))
        except:
            from sqlalchemy import MetaData
            meta = MetaData()
            meta.reflect(bind=eng)
            if "audit_findings" in meta.tables:
                meta.tables["audit_findings"].drop(bind=eng)
                
        Base.metadata.create_all(bind=eng)
        return eng, "Local DB"

engine, db_label = init_db()

def save_findings(uc, findings):
    Session = sessionmaker(bind=engine)
    db = Session()
    # Delete old records for this audit run so Audit Records stays current
    db.query(AuditFinding).filter(AuditFinding.use_case_sl == uc["sl"]).delete()
    uc_name = uc.get("use_case", uc.get("label", "Comprehensive Enterprise Audit"))
    for f in findings:
        db.add(AuditFinding(use_case_sl=uc["sl"], use_case_name=uc_name[:290],
            severity=f.get("severity",""), control=f.get("control",""),
            finding=f.get("finding",""), recommendation=f.get("recommendation",""),
            status=f.get("status","Open"), comment=f.get("comment",""),
            source_files=f.get("source_files","")))
    db.commit(); db.close()

def get_all_findings():
    Session = sessionmaker(bind=engine)
    db = Session()
    rows = db.query(AuditFinding).order_by(AuditFinding.created_at.desc()).all()
    db.close(); return rows

def save_chat_message(session_id, session_title, role, content):
    Session = sessionmaker(bind=engine)
    db = Session()
    db.query(ChatMessage).filter(ChatMessage.session_id == session_id).update({ChatMessage.session_title: session_title})
    db.add(ChatMessage(session_id=session_id, session_title=session_title, role=role, content=content))
    db.commit()
    db.close()

def get_chat_history(session_id):
    Session = sessionmaker(bind=engine)
    db = Session()
    msgs = db.query(ChatMessage).filter(ChatMessage.session_id == session_id).order_by(ChatMessage.created_at.asc()).all()
    db.close()
    return [{"role": m.role, "content": m.content} for m in msgs]

def get_chat_title(session_id):
    Session = sessionmaker(bind=engine)
    db = Session()
    msg = db.query(ChatMessage).filter(ChatMessage.session_id == session_id).order_by(ChatMessage.created_at.asc()).first()
    db.close()
    return msg.session_title if msg else None

def get_all_chat_sessions():
    Session = sessionmaker(bind=engine)
    db = Session()
    from sqlalchemy import func
    sub = db.query(
        ChatMessage.session_id,
        func.max(ChatMessage.created_at).label("last_msg")
    ).group_by(ChatMessage.session_id).subquery()
    rows = db.query(
        ChatMessage.session_id,
        ChatMessage.session_title
    ).join(sub, (ChatMessage.session_id == sub.c.session_id) & (ChatMessage.created_at == sub.c.last_msg)).order_by(sub.c.last_msg.desc()).all()
    db.close()
    
    seen = set()
    sessions = []
    for r in rows:
        if r.session_id not in seen:
            seen.add(r.session_id)
            sessions.append({"session_id": r.session_id, "session_title": r.session_title})
    return sessions

def clear_chat_session(session_id):
    Session = sessionmaker(bind=engine)
    db = Session()
    db.query(ChatMessage).filter(ChatMessage.session_id == session_id).delete()
    db.commit()
    db.close()

def extract_text(f):
    name_lower = f.name.lower()
    if name_lower.endswith(".pdf"):
        with pdfplumber.open(f) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    elif name_lower.endswith((".xlsx", ".xls")):
        try:
            excel_data = pd.read_excel(f, sheet_name=None)
            sheets_text = []
            for sheet_name, df in excel_data.items():
                sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + df.to_string(index=False))
            return "\n\n".join(sheets_text)
        except Exception as e:
            return f"[Error parsing Excel file {f.name}: {e}]"
    elif name_lower.endswith(".csv"):
        try:
            df = pd.read_csv(f)
            return df.to_string(index=False)
        except Exception as e:
            return f"[Error parsing CSV file {f.name}: {e}]"
    elif name_lower.endswith((".pptx", ".ppt")):
        try:
            from pptx import Presentation
            prs = Presentation(f)
            text_runs = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_runs.append(f"--- Slide {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                text_runs.append(" | ".join(row_text))
            return "\n".join(text_runs)
        except Exception as e:
            return f"[Error parsing PowerPoint file {f.name}: {e}]"
    elif name_lower.endswith(".txt"):
        try:
            return f.read().decode("utf-8", errors="ignore")
        except Exception as e:
            return f"[Error parsing text file {f.name}: {e}]"
    else:
        try:
            doc = Document(f)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            return f"[Error parsing Word file {f.name}: {e}]"

def scan_file_security(uploaded_file):
    bytes_data = uploaded_file.getvalue()
    # Reset file pointer for subsequent reading
    uploaded_file.seek(0)
    
    # 1. Offline Signature Check: Prevent executables masquerading as documents
    if bytes_data.startswith(b'MZ'):
        return False, "Executable payload disguised as document (MZ signature detected)."
    
    # 2. Offline Hash Blacklist Check (SHA-256)
    file_hash = hashlib.sha256(bytes_data).hexdigest()
    blacklist = [
        "5e884898da28047151d0e56f8dc6292773603d0d6aabbdd62a11ef721d1542d8" # Demo malicious hash
    ]
    if file_hash in blacklist:
        return False, f"Known malicious hash signature identified ({file_hash[:8]}...)."
        
    return True, "Clean"

def ai_chat_stream(system_ctx, user_msg, model_choice):
    # Professional Auditor Persona
    enhanced_sys = f"You are a Senior Cybersecurity Auditor with expertise in ISO 27001, NIST, and SOC 2. {system_ctx}"
    prompt = f"{enhanced_sys}\n\nUser: {user_msg}\n\nAI Auditor:"
    
    # Map UI selection to Ollama model names
    if "Qwen" in model_choice:
        ollama_model = "qwen2.5:7b"
    else:
        ollama_model = "llama3.1"
        
    try:
        r = requests.post("http://127.0.0.1:11434/api/generate",
            json={"model": ollama_model, "prompt": prompt, "stream": True}, stream=True, timeout=90)
            
        if r.status_code != 200:
            try:
                err = r.json().get("error", r.text)
            except:
                err = r.text
            yield f"⚠️ Ollama Error: {err}. Please make sure you have downloaded the model using pull_models.bat!"
            return

        for line in r.iter_lines():
            if line:
                chunk = json.loads(line)
                yield chunk.get("response", "")
    except Exception as e:
        yield f"⚠️ Offline Engine not responding: {e}"

# ── SESSION STATE ─────────────────────────────────────────────────────────────
if "active_chat_id" not in st.session_state:
    st.session_state.active_chat_id = uuid.uuid4().hex

for k,v in [("stage",0),("context",""),("findings",[]),("chat",[]),("sel_uc",0),("_last_loaded_chat_id","")]:
    if k not in st.session_state: st.session_state[k] = v

# Only reload chat from DB when the session ID actually changes (not on every rerun)
# This is the ChatGPT-style pattern: switching sessions loads history, new chat stays empty
if st.session_state._last_loaded_chat_id != st.session_state.active_chat_id:
    st.session_state.chat = get_chat_history(st.session_state.active_chat_id)
    st.session_state._last_loaded_chat_id = st.session_state.active_chat_id

uc = USE_CASES[st.session_state.sel_uc]

# ── SIDEBAR ───────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("## 🛡️ AICyberAuditBox")
    st.markdown("<small style='color:#64748b'>Agentic RAG Auditor</small>", unsafe_allow_html=True)
    st.markdown(f"<small style='color:#22c55e'>● {db_label} Connected</small>", unsafe_allow_html=True)
    st.divider()

    # ── New Chat button ───────────────────────────────────────────────────────
    if st.button("✏️  New Chat", use_container_width=True, type="primary"):
        new_id = uuid.uuid4().hex
        st.session_state.active_chat_id = new_id
        st.session_state.update({
            "chat": [], "context": "", "findings": [], "stage": 0,
            "resolved_count": None, "resolved_controls": set(),
            "resolved_list": [], "ewaste_resolved": None,
            "last_uploaded_names": "", "_last_loaded_chat_id": new_id
        })
        st.rerun()

    # ── Recents toggle ────────────────────────────────────────────────────────
    sessions = get_all_chat_sessions()
    if sessions:
        if "recents_open" not in st.session_state:
            st.session_state.recents_open = True

        arrow = "▾" if st.session_state.recents_open else "▸"
        if st.button(f"{arrow}  Recents", use_container_width=True, key="recents_toggle"):
            st.session_state.recents_open = not st.session_state.recents_open
            st.rerun()

        if st.session_state.recents_open:
            for s in sessions:
                title = s["session_title"] or "Untitled Chat"
                display = title[:26] + "…" if len(title) > 26 else title
                is_active = s["session_id"] == st.session_state.active_chat_id

                col_t, col_d = st.columns([6, 1])
                with col_t:
                    # Active session highlighted, others plain
                    label = f"**{display}**" if is_active else display
                    if st.button(label, key=f"ch_{s['session_id']}", use_container_width=True,
                                 type="primary" if is_active else "secondary"):
                        st.session_state.active_chat_id = s["session_id"]
                        st.session_state.chat = get_chat_history(s["session_id"])
                        st.session_state._last_loaded_chat_id = s["session_id"]
                        st.rerun()
                with col_d:
                    if st.button("✕", key=f"dx_{s['session_id']}", use_container_width=True):
                        clear_chat_session(s["session_id"])
                        if is_active:
                            new_id = uuid.uuid4().hex
                            st.session_state.active_chat_id = new_id
                            st.session_state.chat = []
                            st.session_state._last_loaded_chat_id = new_id
                        st.rerun()

    st.divider()

    st.markdown("**AI Engine Setup**")
    ai_model = st.selectbox("Select Offline LLM (via Ollama)", [
        "Llama 3.1 (8B) - High Performance Generalist", 
        "Qwen 2.5 (7B) - High Performance Auditor/Reasoning"
    ], label_visibility="collapsed")
    st.divider()

    st.markdown("**Compliance Standard**")
    selected_standard = st.selectbox("Select Target Framework", [
        "All Standards",
        "ISO 27001",
        "DPDP / GDPR",
        "SOC 2",
        "BCMS (Business Continuity)",
        "X-BOM (Software Bill of Materials)"
    ], label_visibility="collapsed")

    # Filter use cases based on standard
    if selected_standard == "All Standards":
        filtered_use_cases = USE_CASES
    else:
        filtered_use_cases = [u for u in USE_CASES if u["standard"] == selected_standard]
       
    st.markdown("**Target Controls to Audit**")
    selected_control_labels = st.multiselect(
        "Select individual controls",
        options=[u["label"] for u in filtered_use_cases],
        default=[u["label"] for u in filtered_use_cases],
        label_visibility="collapsed"
    )
    
    selected_ucs = [u for u in filtered_use_cases if u["label"] in selected_control_labels]
    selected_sls = {u["sl"] for u in selected_ucs}
    st.divider()

    st.markdown("**Upload Evidence**")
    uploaded = st.file_uploader("Upload evidence document(s)", type=["pdf","docx","doc","xlsx","xls","csv","pptx","ppt","txt"],
                                accept_multiple_files=True, label_visibility="collapsed")
    
    # ── REAL-TIME FILE PARSING FOR AI ASSISTANT ──────────────────────────────
    if "last_uploaded_names" not in st.session_state:
        st.session_state.last_uploaded_names = ""
    
    uploaded_names_str = ", ".join([f.name for f in uploaded]) if uploaded else ""
    if (uploaded_names_str != st.session_state.last_uploaded_names) or (uploaded and not st.session_state.context):
        if uploaded:
            auto_ctx = ""
            for f in uploaded:
                try:
                    auto_ctx += f"--- FILE: {f.name} ---\n{extract_text(f)}\n\n"
                except Exception as ex:
                    auto_ctx += f"--- FILE: {f.name} ---\n(Error extracting text: {ex})\n\n"
            st.session_state.context = auto_ctx.strip()
        else:
            st.session_state.context = ""
        st.session_state.last_uploaded_names = uploaded_names_str

    st.divider()

    col_run, col_rst = st.columns([2,1])
    run = col_run.button("▶ Run Analysis", type="primary", use_container_width=True)
    if col_rst.button("↺", use_container_width=True):
        st.session_state.update({"stage":0,"context":"","findings":[],"chat":[],"ewaste_resolved":None})
        clear_chat_session(st.session_state.active_chat_id)
        st.session_state.active_chat_id = uuid.uuid4().hex
        st.rerun()
    
    # Show how many gaps were resolved after analysis
    resolved = st.session_state.get("resolved_count", None)
    if resolved is not None:
        if resolved > 0:
            st.success(f"✅ {resolved} gap(s) resolved by uploaded evidence")
        else:
            st.warning("⚠️ No resolving evidence found in documents")
            with st.expander("🔍 Inspect Extracted Text"):
                if st.session_state.get("context", ""):
                    st.text_area("Extracted Context (First 3000 chars)", st.session_state.context[:3000], height=200, disabled=True)
                else:
                    st.error("No text could be extracted. The document may be empty, password-protected, or a scanned image.")

# ── PIPELINE EXECUTION ────────────────────────────────────────────────────────
if run:
    if not uploaded:
        st.sidebar.error("Please upload the evidence file first.")
    else:
        malware_detected = False
        ctx = ""
        file_texts = {}
        for f in uploaded:
            # 🛡️ Run Offline Malware Scan
            is_clean, reason = scan_file_security(f)
            if not is_clean:
                st.sidebar.error(f"🚨 SECURITY ALERT: '{f.name}' BLOCKED! {reason}")
                malware_detected = True
                break
                
            text = extract_text(f)
            ctx += f"--- FILE: {f.name} ---\n{text}\n\n"
            file_texts[f.name] = text
            
        if malware_detected:
            st.stop()
            
        st.session_state.context = ctx.strip()
        for s in range(1, 5):
            st.session_state.stage = s
            time.sleep(0.3)
        # ── DYNAMIC GAP RESOLUTION ENGINE ──────────────────────────────────────
        # For each finding, check which uploaded documents contain resolving evidence
        resolved_mapping = {} # maps control -> list of file names
        for control, keywords in GAP_RESOLUTION.items():
            matching_files = []
            for fname, ftext in file_texts.items():
                if any(kw in ftext.lower() for kw in keywords):
                    matching_files.append(fname)
            if matching_files:
                resolved_mapping[control] = matching_files

        # Store count for sidebar indicator
        st.session_state["resolved_count"] = len(resolved_mapping)
        st.session_state["resolved_controls"] = set(resolved_mapping.keys())

        all_findings = []
        resolved_list = []
        
        file_names_list = list(file_texts.keys())
        scanned_files_str = ", ".join(file_names_list) if file_names_list else "None"
        
        for k, v in DEMO_FINDINGS.items():
            if k == "CROSS_FILE":
                if len(uploaded) >= 1:
                    for f in v:
                        f_copy = f.copy()
                        f_copy["status"] = "Open"
                        f_copy["comment"] = ""
                        f_copy["editing"] = False
                        
                        # Populate file names in cross-file finding dynamically
                        f_text = f_copy.get("finding", "")
                        if "File 1" in f_text:
                            file1_name = file_names_list[0]
                            file2_name = file_names_list[1] if len(file_names_list) > 1 else file_names_list[0]
                            f_text = f_text.replace("File 1", f"'{file1_name}'").replace("File 2", f"'{file2_name}'")
                            f_copy["finding"] = f_text
                            if len(file_names_list) > 1:
                                f_copy["source_files"] = f"Correlated between: '{file1_name}' and '{file2_name}'"
                            else:
                                f_copy["source_files"] = f"Internal correlation in: '{file1_name}'"
                        else:
                            f_copy["source_files"] = f"Checked in: {scanned_files_str}"
                            
                        all_findings.append(f_copy)
                continue
            if k in selected_sls:
                for finding in v:
                    ctrl = finding.get("control", "")
                    if ctrl in resolved_mapping:
                        # Gap resolved — remove from report, track separately
                        resolved_files = resolved_mapping[ctrl]
                        resolved_list.append(ctrl)
                    else:
                        f_copy = finding.copy()
                        f_copy["status"] = "Open"
                        f_copy["comment"] = ""
                        f_copy["editing"] = False
                        f_copy["source_files"] = f"Checked in: {scanned_files_str} (Evidence missing)"
                        all_findings.append(f_copy)

        st.session_state["resolved_list"] = resolved_list
        st.session_state.findings = all_findings
        st.session_state.stage = 5
        st.rerun()

# ── MAIN LAYOUT ───────────────────────────────────────────────────────────────
# Header
st.markdown(f"""
<div class="main-header">
  <div style="display:flex;align-items:center;gap:16px">
    <div style="font-size:2.5rem">🛡️</div>
    <div>
      <div style="font-size:1.6rem;font-weight:700;color:#f8fafc">AICyberAuditBox</div>
      <div style="color:#64748b;font-size:.9rem">Agentic RAG · Cyber Security Audit Intelligence</div>
    </div>
    <div style="margin-left:auto;text-align:right">
      <div style="color:#22c55e;font-weight:600;font-size:.85rem">● SYSTEM ONLINE</div>
      <div style="color:#64748b;font-size:.8rem">{datetime.now().strftime('%d %b %Y  %H:%M')}</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

# Main column
with st.container():
    tab2, tab1, tab3 = st.tabs(["💬  AI Assistant", "📊  Audit Report", "🗄️  Audit Records"])

    # ── TAB 1 ─────────────────────────────────────────────────────────────────
    with tab1:
        if st.session_state.stage == 0:
            st.markdown("### 📤 Upload Evidence to Begin")
            st.info("Select compliance framework and individual controls in the sidebar, upload your evidence document(s), and click **Run Analysis** to automatically detect security gaps.")

        elif st.session_state.stage == 5:
            findings = st.session_state.findings
            resolved_list = st.session_state.get("resolved_list", [])
            
            # Filter out dismissed/deleted findings for active counts and report view
            active_findings = [f for f in findings if f.get("status", "Open") != "Dismissed"]
            
            counts = {"CRITICAL":0,"HIGH":0,"MEDIUM":0}
            for f in active_findings:
                sev = f.get("severity","MEDIUM").upper()
                if sev in counts: counts[sev] = counts[sev] + 1

            c1,c2,c3,c4 = st.columns(4)
            c1.markdown(f"<div class='stat-card'><div class='stat-num' style='color:#ef4444'>{counts['CRITICAL']}</div><div style='color:#94a3b8'>P1 · Critical</div></div>", unsafe_allow_html=True)
            c2.markdown(f"<div class='stat-card'><div class='stat-num' style='color:#f97316'>{counts['HIGH']}</div><div style='color:#94a3b8'>P2 · High</div></div>", unsafe_allow_html=True)
            c3.markdown(f"<div class='stat-card'><div class='stat-num' style='color:#eab308'>{counts['MEDIUM']}</div><div style='color:#94a3b8'>P3 · Medium</div></div>", unsafe_allow_html=True)
            c4.markdown(f"<div class='stat-card'><div class='stat-num' style='color:#22c55e'>{len(resolved_list)}</div><div style='color:#94a3b8'>✓ Resolved</div></div>", unsafe_allow_html=True)

            # Show resolved controls as a clean green banner
            if resolved_list:
                resolved_html = " &nbsp;·&nbsp; ".join([f"<b>{c}</b>" for c in resolved_list])
                st.markdown(f"<div style='background:rgba(34,197,94,0.1);border:1px solid #22c55e;border-radius:8px;padding:10px 16px;margin:12px 0;color:#22c55e;font-size:0.85rem'>✅ <b>Resolved Controls:</b> &nbsp;{resolved_html}</div>", unsafe_allow_html=True)

            st.markdown(f"<br><small style='color:#64748b'>Generated · {datetime.now().strftime('%d %b %Y %H:%M:%S')} · {selected_standard} ({len(selected_ucs)} Controls)</small>", unsafe_allow_html=True)
            st.divider()

            SEVERITY_LABEL = {"CRITICAL": "P1 · CRITICAL", "HIGH": "P2 · HIGH", "MEDIUM": "P3 · MEDIUM"}
            CSS = {"CRITICAL":"badge-critical","HIGH":"badge-high","MEDIUM":"badge-medium"}
            EMJ = {"CRITICAL":"🔴","HIGH":"🟠","MEDIUM":"🟡"}
            
            # Sort findings by severity
            open_findings_sorted = sorted(active_findings, key=lambda x: ["CRITICAL","HIGH","MEDIUM"].index(x.get("severity","MEDIUM").upper()))
            
            for idx, f in enumerate(open_findings_sorted):
                s = f.get("severity","MEDIUM").upper()
                label = SEVERITY_LABEL.get(s, s)
                css = CSS.get(s, "badge-medium")
                emj = EMJ.get(s, "🟡")
                status = f.get("status", "Open")
                editing = f.get("editing", False)
                
                status_color = "#3b82f6" if status == "Open" else "#22c55e"
                
                if editing:
                    # Render inline editor container
                    with st.container(border=True):
                        st.markdown("##### ✏️ Modify Finding Details")
                        col_edit_sev, col_edit_ctrl = st.columns([1, 2])
                        with col_edit_sev:
                            sev_index = ["CRITICAL", "HIGH", "MEDIUM"].index(s) if s in ["CRITICAL", "HIGH", "MEDIUM"] else 2
                            new_sev = st.selectbox("Severity", ["CRITICAL", "HIGH", "MEDIUM"], index=sev_index, key=f"sev_edit_sel_{idx}")
                        with col_edit_ctrl:
                            new_ctrl = st.text_input("Control", value=f.get("control", ""), key=f"ctrl_edit_in_{idx}")
                        
                        new_finding = st.text_area("Finding Description", value=f.get("finding", ""), key=f"find_edit_ta_{idx}", height=80)
                        new_rec = st.text_area("Recommendation/Mitigation", value=f.get("recommendation", ""), key=f"rec_edit_ta_{idx}", height=80)
                        new_src = st.text_input("Source File Scope", value=f.get("source_files", "All uploaded documents"), key=f"src_edit_in_{idx}")
                        
                        col_save, col_cancel = st.columns([1.5, 1.5])
                        with col_save:
                            if st.button("💾 Save Changes", key=f"save_edit_{idx}", type="primary", use_container_width=True):
                                for orig_f in st.session_state.findings:
                                    if orig_f["control"] == f["control"] and orig_f["finding"] == f["finding"]:
                                        orig_f["severity"] = new_sev
                                        orig_f["control"] = new_ctrl
                                        orig_f["finding"] = new_finding
                                        orig_f["recommendation"] = new_rec
                                        orig_f["source_files"] = new_src
                                        orig_f["editing"] = False
                                st.rerun()
                        with col_cancel:
                            if st.button("Cancel", key=f"cancel_edit_{idx}", use_container_width=True):
                                for orig_f in st.session_state.findings:
                                    if orig_f["control"] == f["control"] and orig_f["finding"] == f["finding"]:
                                        orig_f["editing"] = False
                                st.rerun()
                else:
                    # Render standard static view
                    st.markdown(f"""
                    <div class='{css}' style='margin-bottom:0px; border-bottom-left-radius:0px; border-bottom-right-radius:0px;'>
                      <div style='display:flex; justify-content:space-between; align-items:center;'>
                        <b>{emj} {label}</b>
                        <span style='font-size:0.75rem; background:{status_color}; color:white; padding:2px 8px; border-radius:12px; font-weight:600;'>{status.upper()}</span>
                      </div>
                      <div style='margin-top:6px;'><b>Control:</b> {f.get('control','')}</div>
                      <span style='color:#cbd5e1'>📌 <b>Finding:</b> {f.get('finding','')}</span><br>
                      <span style='color:#86efac'>→ <b>Recommendation:</b> {f.get('recommendation','')}</span>
                      <div style='margin-top:8px; font-size:0.8rem; color:#94a3b8; border-top:1px dashed #334155; padding-top:6px; display:flex; align-items:center; gap:6px;'>
                        <span>📁</span> <b>Source File Scope:</b> <i>{f.get('source_files','All uploaded documents')}</i>
                      </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Attached actions section
                    with st.container(border=True):
                        col_act1, col_act2, col_act3, col_act4 = st.columns([1.8, 1.8, 1.8, 5])
                        with col_act1:
                            if status == "Accepted":
                                if st.button("↩ Undo", key=f"undo_{idx}", use_container_width=True, type="secondary"):
                                    for orig_f in st.session_state.findings:
                                        if orig_f["control"] == f["control"] and orig_f["finding"] == f["finding"]:
                                            orig_f["status"] = "Open"
                                    st.rerun()
                            else:
                                if st.button("✓ Accept", key=f"acc_{idx}", use_container_width=True, type="secondary"):
                                    for orig_f in st.session_state.findings:
                                        if orig_f["control"] == f["control"] and orig_f["finding"] == f["finding"]:
                                            orig_f["status"] = "Accepted"
                                    st.rerun()
                        with col_act2:
                            if st.button("✏️ Modify", key=f"mod_{idx}", use_container_width=True, type="secondary"):
                                for orig_f in st.session_state.findings:
                                    if orig_f["control"] == f["control"] and orig_f["finding"] == f["finding"]:
                                        orig_f["editing"] = True
                                st.rerun()
                        with col_act3:
                            if st.button("🗑️ Delete", key=f"del_{idx}", use_container_width=True, type="secondary"):
                                for orig_f in st.session_state.findings:
                                    if orig_f["control"] == f["control"] and orig_f["finding"] == f["finding"]:
                                        orig_f["status"] = "Dismissed"
                                st.rerun()
                        with col_act4:
                            comment_val = st.text_input("Auditor Notes", value=f.get("comment", ""), key=f"cmt_{idx}", label_visibility="collapsed", placeholder="Add auditor notes or comments...")
                            if comment_val != f.get("comment", ""):
                                for orig_f in st.session_state.findings:
                                    if orig_f["control"] == f["control"] and orig_f["finding"] == f["finding"]:
                                        orig_f["comment"] = comment_val


            # Display restore option for deleted findings
            dismissed_findings = [df for df in findings if df.get("status", "Open") == "Dismissed"]
            if dismissed_findings:
                st.markdown("<br>", unsafe_allow_html=True)
                with st.expander(f"🗑️ Deleted Findings ({len(dismissed_findings)})", expanded=False):
                    for idx_d, df in enumerate(dismissed_findings):
                        col_text, col_restore = st.columns([8, 2])
                        with col_text:
                            st.markdown(f"**{df.get('control', '')}** — <span style='color:#94a3b8'>{df.get('finding', '')[:90]}...</span>", unsafe_allow_html=True)
                        with col_restore:
                            if st.button("↩ Restore", key=f"restore_{idx_d}", use_container_width=True):
                                for orig_f in st.session_state.findings:
                                    if orig_f["control"] == df["control"] and orig_f["finding"] == df["finding"]:
                                        orig_f["status"] = "Open"
                                st.rerun()

            st.divider()
            b1, b2 = st.columns(2)
            with b1:
                if st.button("💾  Save to ShaktiDB", type="primary", use_container_width=True):
                    save_findings({"sl": 0, "use_case": f"{selected_standard} Audit Run"}, active_findings)
                    st.success(f"✅ {len(active_findings)} findings saved to {db_label}")
            with b2:
                df_export = pd.DataFrame([{
                    "Control": f.get("control", ""),
                    "Severity": f.get("severity", ""),
                    "Finding": f.get("finding", ""),
                    "Recommendation": f.get("recommendation", ""),
                    "Status": f.get("status", "Open"),
                    "Source Scope": f.get("source_files", "All uploaded documents"),
                    "Auditor Comment": f.get("comment", "")
                } for f in active_findings])
                csv_data = df_export.to_csv(index=False)
                st.download_button("⬇️  Export Report CSV", csv_data, "comprehensive_audit_report.csv", use_container_width=True)

    # ── TAB 2 ─────────────────────────────────────────────────────────────────
    with tab2:
        st.markdown("""
        <div style='background:#1e293b;border:1px solid #334155;border-radius:12px;padding:16px;margin-bottom:16px;display:flex;align-items:center;gap:12px'>
          <div style='font-size:2rem'>🤖</div>
          <div>
            <div style='font-weight:700;color:#f8fafc'>AI Audit Assistant</div>
            <div style='color:#64748b;font-size:.85rem'>Local LLM · No internet required · Evidence-aware</div>
          </div>
          <div style='margin-left:auto;color:#22c55e;font-size:.8rem;font-weight:600'>● ONLINE</div>
        </div>""", unsafe_allow_html=True)
        
        if len(st.session_state.context) > 0:
            st.markdown("<div style='background:rgba(59,130,246,0.1); border:1px solid #3b82f6; border-radius:8px; padding:8px 12px; color:#3b82f6; font-size:0.85rem; font-weight:600; margin-bottom:16px'>🔍 Cross-File Intelligence Active · Correlating multiple evidence sources</div>", unsafe_allow_html=True)

        if st.session_state.context:
            st.success(f"✅ Evidence document loaded · {len(st.session_state.context):,} characters indexed")
        else:
            st.info("💡 Upload and run analysis first for evidence-aware answers, or ask general cybersecurity questions.")

        for msg in st.session_state.chat:
            if msg["role"] == "user":
                st.markdown(f"<div style='text-align:right;font-size:11px;color:#64748b;margin-top:8px'>You</div><div class='chat-bubble-user'>{msg['content']}</div>", unsafe_allow_html=True)
            else:
                st.markdown(f"<div style='font-size:11px;color:#3b82f6;font-weight:600;margin-top:8px'>🤖 AI Auditor</div><div class='chat-bubble-bot'>{msg['content']}</div>", unsafe_allow_html=True)

        user_msg = st.chat_input("Ask the AI Auditor anything...")
        if user_msg:
            # Determine active title
            title = get_chat_title(st.session_state.active_chat_id)
            if not title:
                title = user_msg[:30] + ("..." if len(user_msg) > 30 else "")
            
            save_chat_message(st.session_state.active_chat_id, title, "user", user_msg)
            
            # Immediately show the user's message in the UI so it doesn't freeze
            st.markdown(f"<div style='text-align:right;font-size:11px;color:#64748b;margin-top:8px'>You</div><div class='chat-bubble-user'>{user_msg}</div>", unsafe_allow_html=True)
            
            sys = "You are a Senior Cybersecurity Auditor. PERFORM CROSS-DOCUMENT CORRELATION: Look for inconsistencies, contradictions, or missing links between the multiple uploaded files. If File A mentions a policy but File B shows it is not followed, flag it. Be precise, professional, and structured."
            if st.session_state.context:
                sys += f"\n\nEVIDENCE:\n{st.session_state.context[:4000]}"
            if st.session_state.findings:
                sys += f"\n\nFINDINGS:\n{json.dumps(st.session_state.findings)[:1500]}"
            
            # Setup real-time streaming UI
            st.markdown(f"<div style='font-size:11px;color:#3b82f6;font-weight:600;margin-top:8px'>🤖 AI Auditor ({ai_model.split(' ')[0]})</div>", unsafe_allow_html=True)
            placeholder = st.empty()
            full_ans = ""
            
            # Stream the conversation tokens as they arrive
            for chunk in ai_chat_stream(sys, user_msg, ai_model):
                full_ans += chunk
                placeholder.markdown(f"<div class='chat-bubble-bot'>{full_ans}▌</div>", unsafe_allow_html=True)
            
            # Final output without the cursor block
            placeholder.markdown(f"<div class='chat-bubble-bot'>{full_ans}</div>", unsafe_allow_html=True)
            
            save_chat_message(st.session_state.active_chat_id, title, "assistant", full_ans)
            st.rerun()

        if st.session_state.chat:
            if st.button("🗑️ Clear Active Chat", use_container_width=True):
                clear_chat_session(st.session_state.active_chat_id)
                st.rerun()

    # ── TAB 3 ─────────────────────────────────────────────────────────────────
    with tab3:
        st.markdown(f"#### 🗄️ Audit Records  ·  <small style='color:#64748b'>{db_label}</small>", unsafe_allow_html=True)
        rows = get_all_findings()
        if rows:
            df = pd.DataFrame([{
                "UC": f"UC{r.use_case_sl}",
                "Scenario": (r.use_case_name or "")[:55],
                "Severity": r.severity,
                "Control": r.control,
                "Finding": (r.finding or "")[:90],
                "Recommendation": (r.recommendation or "")[:90],
                "Status": r.status,
                "Source Scope": r.source_files,
                "Comment": r.comment,
                "Date": r.created_at.strftime("%d %b %Y") if r.created_at else ""
            } for r in rows])
            st.dataframe(df, use_container_width=True, hide_index=True)
            col_exp, col_clear = st.columns(2)
            with col_exp:
                st.download_button("⬇️ Export All Records", df.to_csv(index=False), "all_audit_findings.csv", use_container_width=True)
            with col_clear:
                if st.button("🗑️ Clear All Database Records", use_container_width=True, type="secondary"):
                    Session = sessionmaker(bind=engine)
                    db = Session()
                    db.query(AuditFinding).delete()
                    db.commit()
                    db.close()
                    st.success("✅ Database records cleared successfully!")
                    st.rerun()
        else:
            st.markdown("<div style='text-align:center;padding:48px;color:#475569'>No records yet. Run an audit and save findings.</div>", unsafe_allow_html=True)

st.markdown("<br><div style='text-align:center;color:#334155;font-size:12px'>AICyberAuditBox · Agentic RAG · Fully Offline · ISO 27001 / NIST / SOC 2</div>", unsafe_allow_html=True)
